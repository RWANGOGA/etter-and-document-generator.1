import io
import json
import re
from fastapi import APIRouter, HTTPException, UploadFile, File, Form
from fastapi.responses import Response
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from groq import Groq

from app.config import settings

router = APIRouter(prefix="/api/pptx", tags=["pptx"])
client = Groq(api_key=settings.groq_api_key)


class ThemeOptions(BaseModel):
    primary_color: str = "1E3A5F"
    accent_color: str = "2B6CB0"
    text_color: str = "333333"
    font_family: str = "Calibri"


class PptxRequest(BaseModel):
    topic: str
    slide_count: int = 8
    theme: ThemeOptions = ThemeOptions()


class ChatMessage(BaseModel):
    role: str
    content: str


class ChatRequest(BaseModel):
    messages: list[ChatMessage]


class ChatResponse(BaseModel):
    success: bool
    reply: str
    structure: dict | None = None
    ready: bool = False


class BuildRequest(BaseModel):
    structure: dict
    theme: ThemeOptions = ThemeOptions()


def _hex_to_rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str.lstrip("#").upper())


def _strip_thinking(text: str) -> str:
    return re.sub(r"<think>.*?</think>", "", text, flags=re.DOTALL).strip()


def _extract_json(text: str) -> dict:
    text = _strip_thinking(text)
    text = re.sub(r"^```(?:json)?\s*|\s*```$", "", text.strip())
    start = text.find("{")
    end = text.rfind("}") + 1
    if start == -1 or end == 0:
        raise ValueError("No JSON object found in model output")
    return json.loads(text[start:end])


STRUCTURE_PROMPT = """You are generating content for a professional PowerPoint presentation.
Topic: {topic}
Target slide count: {slide_count}

Return ONLY valid JSON, no markdown fences, no commentary, in exactly this shape:
{{
  "title": "Presentation title",
  "subtitle": "One-line subtitle",
  "slides": [
    {{"heading": "Slide heading", "bullets": ["point one", "point two", "point three"]}}
  ]
}}

Rules:
- First slide is the title slide (handled separately, do not include it in "slides").
- Each slide needs 3-5 concise, substantive bullets — not generic filler.
- Cover the topic with real depth: context, specifics, examples, implications.
- Keep each bullet under 20 words.
- Produce exactly {slide_count} content slides in the "slides" array."""


def _generate_structure(topic: str, slide_count: int) -> dict:
    prompt = STRUCTURE_PROMPT.format(topic=topic, slide_count=slide_count)
    completion = client.chat.completions.create(
        model="qwen/qwen3.6-27b",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=4096,
        temperature=0.6,
    )
    raw = completion.choices[0].message.content or ""
    return _extract_json(raw)


FROM_DOC_PROMPT = """You are turning an existing document into a PowerPoint presentation.
Target slide count: {slide_count}
Extra instructions from the user (may be empty): {instructions}

Source document content:
{content}

Return ONLY valid JSON, no markdown fences, no commentary, in exactly this shape:
{{
  "title": "Presentation title",
  "subtitle": "One-line subtitle",
  "slides": [
    {{"heading": "Slide heading", "bullets": ["point one", "point two", "point three"]}}
  ]
}}

Rules:
- Distill the document's actual content into slides — don't invent unrelated material.
- First slide is the title slide (handled separately, do not include it in "slides").
- Each slide needs 3-5 concise bullets summarizing that section of the source.
- Keep each bullet under 20 words.
- Produce exactly {slide_count} content slides."""


def _extract_text_from_upload(filename: str, content: bytes) -> str:
    name = (filename or "").lower()
    if name.endswith(".pdf"):
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(content))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    elif name.endswith(".docx"):
        import docx
        d = docx.Document(io.BytesIO(content))
        return "\n".join(p.text for p in d.paragraphs)
    else:
        return content.decode("utf-8", errors="ignore")


def _generate_from_document(text: str, slide_count: int, instructions: str) -> dict:
    prompt = FROM_DOC_PROMPT.format(
        slide_count=slide_count,
        instructions=instructions or "None",
        content=text[:12000],
    )
    completion = client.chat.completions.create(
        model="qwen/qwen3.6-27b",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=4096,
        temperature=0.5,
    )
    raw = completion.choices[0].message.content or ""
    return _extract_json(raw)


CHAT_SYSTEM_PROMPT = """You are a presentation-building assistant. Have a short conversation with the user to understand what presentation they need: its purpose, audience, key points to cover, and roughly how many slides.

Ask at most 2-3 clarifying questions total — don't drag it out. Once you have enough to build something genuinely useful (or the user says to just go ahead), respond with:
1. One short plain-text confirmation line (no markdown, no emojis).
2. Then the full structure as JSON inside <STRUCTURE>...</STRUCTURE> tags, in exactly this shape:
{
  "title": "Presentation title",
  "subtitle": "One-line subtitle",
  "slides": [{"heading": "Slide heading", "bullets": ["point one", "point two"]}]
}

Do not include <STRUCTURE> tags until ready to deliver the final structure — during clarifying questions, talk normally in plain text, no markdown symbols."""


def _style_title_slide(slide, title: str, subtitle: str, theme: ThemeOptions):
    primary = _hex_to_rgb(theme.primary_color)
    accent = _hex_to_rgb(theme.accent_color)

    slide.shapes.title.text = title
    p = slide.shapes.title.text_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.color.rgb = primary
    p.font.bold = True
    p.font.name = theme.font_family

    if len(slide.placeholders) > 1:
        sub = slide.placeholders[1]
        sub.text = subtitle
        sp = sub.text_frame.paragraphs[0]
        sp.font.size = Pt(20)
        sp.font.color.rgb = accent
        sp.font.name = theme.font_family


def _style_content_slide(slide, heading: str, bullets: list[str], theme: ThemeOptions):
    primary = _hex_to_rgb(theme.primary_color)
    text_color = _hex_to_rgb(theme.text_color)

    slide.shapes.title.text = heading
    tp = slide.shapes.title.text_frame.paragraphs[0]
    tp.font.size = Pt(30)
    tp.font.color.rgb = primary
    tp.font.bold = True
    tp.font.name = theme.font_family

    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = text_color
        p.font.name = theme.font_family
        p.space_after = Pt(10)


def build_pptx(structure: dict, theme: ThemeOptions) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    title_slide = prs.slides.add_slide(title_layout)
    _style_title_slide(title_slide, structure.get("title", "Untitled"), structure.get("subtitle", ""), theme)

    for slide_data in structure.get("slides", []):
        slide = prs.slides.add_slide(content_layout)
        _style_content_slide(slide, slide_data.get("heading", ""), slide_data.get("bullets", []), theme)

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


@router.post("/generate")
async def generate_pptx(req: PptxRequest):
    if not req.topic.strip():
        raise HTTPException(status_code=400, detail="Topic is required")
    try:
        structure = _generate_structure(req.topic, req.slide_count)
        data = build_pptx(structure, req.theme)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Presentation generation failed: {e}")
    return Response(
        content=data,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=presentation.pptx"},
    )


@router.post("/from-document")
async def pptx_from_document(
    file: UploadFile = File(...),
    slide_count: int = Form(8),
    instructions: str = Form(""),
    primary_color: str = Form("1E3A5F"),
    accent_color: str = Form("2B6CB0"),
    text_color: str = Form("333333"),
    font_family: str = Form("Calibri"),
):
    content = await file.read()
    try:
        text = _extract_text_from_upload(file.filename or "", content)
        if not text.strip():
            raise HTTPException(status_code=400, detail="Could not extract any text from the uploaded file")
        structure = _generate_from_document(text, slide_count, instructions)
        theme = ThemeOptions(
            primary_color=primary_color,
            accent_color=accent_color,
            text_color=text_color,
            font_family=font_family,
        )
        data = build_pptx(structure, theme)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Presentation generation failed: {e}")
    return Response(
        content=data,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=presentation.pptx"},
    )


@router.post("/chat", response_model=ChatResponse)
async def chat_pptx(req: ChatRequest):
    messages = [{"role": "system", "content": CHAT_SYSTEM_PROMPT}]
    for m in req.messages:
        messages.append({"role": m.role, "content": m.content})

    try:
        completion = client.chat.completions.create(
            model="qwen/qwen3.6-27b",
            messages=messages,
            max_tokens=4096,
            temperature=0.6,
        )
        text = _strip_thinking(completion.choices[0].message.content or "")

        structure = None
        reply_text = text
        ready = False

        if "<STRUCTURE>" in text and "</STRUCTURE>" in text:
            start = text.index("<STRUCTURE>")
            end = text.index("</STRUCTURE>") + len("</STRUCTURE>")
            raw_json = text[start + len("<STRUCTURE>"):text.index("</STRUCTURE>")].strip()
            reply_text = (text[:start] + text[end:]).strip()
            structure = json.loads(raw_json)
            ready = True
            if not reply_text:
                reply_text = "Here's the presentation — review it and click Build to download."

        return ChatResponse(success=True, reply=reply_text, structure=structure, ready=ready)
    except Exception as e:
        return ChatResponse(success=False, reply=f"Something went wrong: {e}")


@router.post("/build")
async def build_from_structure(req: BuildRequest):
    try:
        data = build_pptx(req.structure, req.theme)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Build failed: {e}")
    return Response(
        content=data,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=presentation.pptx"},
    )